"""Встроенные изображения в .pptx — тот же архитектурный пробел, что был у
.docx (текст слайдов извлекался, картинки — нет), закрытый той же техникой.
Студенты сдают презентации со скриншотами/фото/диаграммами на слайдах не
реже, чем документы с картинкой вместо текста."""
import asyncio
import io

import pytest
from pptx import Presentation
from pptx.util import Inches
from PIL import Image

from services import ai_grader


def _png_bytes(color: tuple) -> bytes:
    buf = io.BytesIO()
    Image.new("RGB", (20, 20), color=color).save(buf, format="PNG")
    return buf.getvalue()


def _build_pptx(images_per_slide: list) -> bytes:
    """images_per_slide: список списков байт картинок, один список на слайд."""
    prs = Presentation()
    blank_layout = prs.slide_layouts[6]
    for slide_images in images_per_slide:
        slide = prs.slides.add_slide(blank_layout)
        for i, img in enumerate(slide_images):
            slide.shapes.add_picture(io.BytesIO(img), Inches(i), Inches(0), height=Inches(1))
    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


IMAGE_1 = _png_bytes((255, 0, 0))
IMAGE_2 = _png_bytes((0, 255, 0))
IMAGE_3 = _png_bytes((0, 0, 255))


def test_extract_pptx_images_finds_all_slides():
    data = _build_pptx([[IMAGE_1], [IMAGE_2]])
    images = ai_grader._extract_pptx_images(data)
    raw = [d for d, _ in images]
    assert len(images) == 2
    assert IMAGE_1 in raw and IMAGE_2 in raw


def test_extract_pptx_images_ordered_by_slide_number_not_lexicographic():
    """slide10.xml < slide2.xml лексикографически, но слайд 2 идёт раньше
    слайда 10 в самой презентации — сортировка обязана быть числовой."""
    slides = [[_png_bytes((i, i, i))] for i in range(1, 13)]  # 12 слайдов
    data = _build_pptx(slides)
    images = ai_grader._extract_pptx_images(data, max_images=20)
    raw = [d for d, _ in images]
    assert len(images) == 12
    # Порядок извлечения должен совпадать с порядком добавления слайдов.
    expected_order = [img for (img,) in slides]
    assert raw == expected_order


def test_extract_pptx_images_deduplicates_identical_image_across_slides():
    data = _build_pptx([[IMAGE_1], [IMAGE_1]])  # одна и та же картинка на 2 слайдах
    images = ai_grader._extract_pptx_images(data)
    raw = [d for d, _ in images]
    assert raw.count(IMAGE_1) == 1


def test_extract_pptx_images_respects_max_images_cap():
    slides = [[_png_bytes((i, 0, 0))] for i in range(ai_grader.MAX_DOCX_EMBEDDED_IMAGES + 3)]
    data = _build_pptx(slides)
    images = ai_grader._extract_pptx_images(data)
    assert len(images) == ai_grader.MAX_DOCX_EMBEDDED_IMAGES


def test_extract_pptx_images_empty_for_pptx_without_images():
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    tb = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(2), Inches(1))
    tb.text_frame.text = "Просто текст"
    buf = io.BytesIO()
    prs.save(buf)
    assert ai_grader._extract_pptx_images(buf.getvalue()) == []


def test_extract_pptx_images_corrupted_bytes_returns_empty_not_raise():
    assert ai_grader._extract_pptx_images(b"not a pptx at all") == []


# ── диспетчер _fetch_embedded_images: docx/pptx/прочее ──────────────────────

class _FakeGetResp:
    def __init__(self, content: bytes):
        self.content = content
        self.is_success = True
        self.status_code = 200
        self.headers = {"content-type": "application/octet-stream"}


class _FakeGetClient:
    def __init__(self, content: bytes):
        self._content = content

    def __call__(self, *a, **k):
        return self

    async def __aenter__(self):
        return self

    async def __aexit__(self, *a):
        return False

    async def get(self, *a, **k):
        return _FakeGetResp(self._content)


def test_fetch_embedded_images_dispatches_pptx(monkeypatch):
    data = _build_pptx([[IMAGE_1]])
    monkeypatch.setattr(ai_grader.httpx, "AsyncClient", _FakeGetClient(data))
    uris = asyncio.run(
        ai_grader._fetch_embedded_images("http://localhost:8000/api/uploads/slides.pptx")
    )
    assert len(uris) == 1
    assert uris[0].startswith("data:image/png;base64,")


def test_fetch_embedded_images_empty_for_pdf(monkeypatch):
    data = _build_pptx([[IMAGE_1]])
    monkeypatch.setattr(ai_grader.httpx, "AsyncClient", _FakeGetClient(data))
    uris = asyncio.run(
        ai_grader._fetch_embedded_images("http://localhost:8000/api/uploads/hw.pdf")
    )
    assert uris == []
